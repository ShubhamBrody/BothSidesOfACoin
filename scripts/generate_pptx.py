"""Generate Ideation Phase PowerPoint for BothSidesOfACoin Hackathon."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Slate 900
SURFACE = RGBColor(0x1E, 0x29, 0x3B)        # Slate 800
CARD_BG = RGBColor(0x33, 0x41, 0x55)        # Slate 700
PRIMARY = RGBColor(0x25, 0x63, 0xEB)        # Royal Blue
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)    # Left Blue
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)     # Right Red
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)         # Neutral Purple
GREEN = RGBColor(0x10, 0xB9, 0x81)          # Success Green
AMBER = RGBColor(0xF5, 0x9E, 0x0B)          # Warning Amber
WHITE = RGBColor(0xF8, 0xFA, 0xFC)          # Slate 50
GRAY = RGBColor(0x94, 0xA3, 0xB8)           # Slate 400
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)     # Slate 300

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, spacing=Pt(8)):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txbox


# ═══════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide, BG_DARK)

# Accent bar at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), PRIMARY)

# Coin icon placeholder
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY
circle.line.fill.background()
add_text(slide, Inches(6.05), Inches(1.55), Inches(1.2), Inches(0.8), "⚖️",
         font_size=44, alignment=PP_ALIGN.CENTER)

# Title
add_text(slide, Inches(2.5), Inches(3.0), Inches(8.3), Inches(1.2),
         "BothSidesOfACoin", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
add_text(slide, Inches(2.5), Inches(4.1), Inches(8.3), Inches(0.7),
         "We don't tell you what to think — we help you think.",
         font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.6),
         "AI-Powered Balanced News & Opinion Platform",
         font_size=18, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

# Theme badge
badge = add_shape(slide, Inches(4.2), Inches(5.8), Inches(4.9), Inches(0.55), SURFACE, corner_radius=0.15)
add_text(slide, Inches(4.3), Inches(5.85), Inches(4.7), Inches(0.45),
         "🔥 Hackathon Theme 1: Reimagining Learning & Collaboration",
         font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)

# Footer
add_text(slide, Inches(3.5), Inches(6.7), Inches(6.3), Inches(0.4),
         "Ideation Phase  •  April 2026  •  Team BothSidesOfACoin",
         font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), RED_ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "🚨 THE PROBLEM", font_size=32, color=RED_ACCENT, bold=True)

# Big quote
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
         '"People don\'t lack information — they lack balanced information."',
         font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Problem cards
problems = [
    ("🔄", "Echo Chambers", "Algorithms feed you what you\nalready believe, creating\nideological isolation."),
    ("📱", "Engagement Over Truth", "Social media optimizes for\nclicks and outrage, not\naccuracy or balance."),
    ("🙈", "No Opposing Views", "Users rarely encounter\nperspectives that challenge\ntheir existing beliefs."),
    ("📉", "Polarization Crisis", "Result: misinformation,\nshallow understanding,\nand societal division."),
]

for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.6 + i * 3.1)
    y = Inches(2.8)
    card = add_shape(slide, x, y, Inches(2.8), Inches(3.6), SURFACE, corner_radius=0.05)

    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.7),
             icon, font_size=36, alignment=PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.2), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.6), Inches(2.2), Inches(1.6),
             desc, font_size=13, color=GRAY)

# Stats bar
add_shape(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), CARD_BG, corner_radius=0.1)
add_text(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.5),
         "📊 64% of Americans say social media has a mostly negative effect on the country  •  "
         "77% see the other party as immoral  •  Misinformation costs the global economy $78B/year",
         font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 3: OUR SOLUTION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "💡 OUR SOLUTION", font_size=32, color=GREEN, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.7),
         "An AI-powered platform that shows ALL sides of every story — breaking echo chambers through balanced perspectives.",
         font_size=18, color=LIGHT_GRAY)

# Solution description
add_text(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
         "How It Works:", font_size=22, color=WHITE, bold=True)

steps = [
    "1️⃣  Collects news from multiple diverse sources",
    "2️⃣  AI classifies each article as Left, Right, or Neutral",
    "3️⃣  Generates multi-perspective summaries side-by-side",
    "4️⃣  Tracks your reading bias and recommends balance",
    "5️⃣  AI Debate Mode lets you explore all viewpoints interactively",
]
add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), steps, font_size=16, color=LIGHT_GRAY)

# 3-Side View mockup
add_shape(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.8), SURFACE, corner_radius=0.05)
add_text(slide, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.4),
         "3-Side View  (Main USP)", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Left column
add_shape(slide, Inches(7.2), Inches(2.6), Inches(1.6), Inches(3.8), RGBColor(0x1E, 0x3A, 0x5F), corner_radius=0.03)
add_text(slide, Inches(7.3), Inches(2.7), Inches(1.4), Inches(0.4),
         "🟦 LEFT", font_size=13, color=BLUE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.3), Inches(3.2), Inches(1.4), Inches(2.8),
         "Progressive groups\nargue that strong\nregulation is\nessential to protect\ncitizens from AI\nrisks and deepen\nequality...",
         font_size=10, color=GRAY)

# Neutral column
add_shape(slide, Inches(8.95), Inches(2.6), Inches(1.6), Inches(3.8), RGBColor(0x2D, 0x21, 0x4C), corner_radius=0.03)
add_text(slide, Inches(9.05), Inches(2.7), Inches(1.4), Inches(0.4),
         "⚪ NEUTRAL", font_size=13, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.05), Inches(3.2), Inches(1.4), Inches(2.8),
         "The EU AI Act was\nproposed in April\n2021 and officially\nadopted in March\n2024. It applies\na risk-based\napproach...",
         font_size=10, color=GRAY)

# Right column
add_shape(slide, Inches(10.7), Inches(2.6), Inches(1.6), Inches(3.8), RGBColor(0x5F, 0x1E, 0x1E), corner_radius=0.03)
add_text(slide, Inches(10.8), Inches(2.7), Inches(1.4), Inches(0.4),
         "🟥 RIGHT", font_size=13, color=RED_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(10.8), Inches(3.2), Inches(1.4), Inches(2.8),
         "Free market\nadvocates warn\nthat excessive\nregulation will\ndrive innovation\noffshore and harm\ncompetitiveness...",
         font_size=10, color=GRAY)

# Tagline
add_shape(slide, Inches(3.5), Inches(6.6), Inches(6.3), Inches(0.55), CARD_BG, corner_radius=0.1)
add_text(slide, Inches(3.6), Inches(6.65), Inches(6.1), Inches(0.45),
         "\"We are not building another news app. We are building a system that helps people think, not just react.\"",
         font_size=13, color=AMBER, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 4: CORE FEATURES
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "🧩 CORE FEATURES", font_size=32, color=PRIMARY, bold=True)

features = [
    ("🔀", "3-Side View", "Every topic shown from Left,\nRight, and Neutral perspectives\nside-by-side with source citations.", PRIMARY),
    ("🧠", "Bias Detection Engine", "Tracks reading patterns, computes\na personal Bias Score (0-100),\nalerts when you're in an echo chamber.", AMBER),
    ("🗣️", "AI Debate Mode", "Chat with Left AI, Right AI, and\nNeutral AI. Watch them debate.\nJudges will LOVE this. 🔥", RED_ACCENT),
    ("📊", "Opinion Balancer", "Dashboard showing your ideology\ndistribution. Gamification with\nbadges for balanced readers.", GREEN),
    ("📅", "Timeline View", "How a conflict started, key events,\nstakeholders, cause-effect chains.\n\"Why is this happening?\"", PURPLE),
    ("🌍", "Impact Analyzer", "\"How does this affect YOU?\"\nEconomic, social, political impact\nanalysis personalized to the user.", BLUE_ACCENT),
]

for i, (icon, title, desc, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.3 + row * 3.0)

    card = add_shape(slide, x, y, Inches(3.8), Inches(2.7), SURFACE, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.6),
             icon, font_size=32)
    add_text(slide, x + Inches(1.0), y + Inches(0.25), Inches(2.5), Inches(0.4),
             title, font_size=18, color=accent, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.5),
             desc, font_size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════
# SLIDE 5: AI ARCHITECTURE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), PURPLE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "🧠 AI ARCHITECTURE — Microsoft AutoGen 0.4", font_size=28, color=PURPLE, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.5),
         "Dynamic multi-agent orchestration — NO hardcoded pipelines. The AI decides what to do next.",
         font_size=16, color=LIGHT_GRAY)

# Orchestrator in center
orch = add_shape(slide, Inches(5.2), Inches(1.9), Inches(2.9), Inches(1.2), PRIMARY, corner_radius=0.05)
add_text(slide, Inches(5.3), Inches(2.0), Inches(2.7), Inches(0.5),
         "🧠 Orchestrator Agent", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.3), Inches(2.5), Inches(2.7), Inches(0.5),
         "Plans steps dynamically\nusing LLM reasoning", font_size=11, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)

# Agent cards
agents = [
    ("📰", "News\nCollector", Inches(0.6), Inches(3.6)),
    ("⚖️", "Bias\nClassifier", Inches(2.8), Inches(3.6)),
    ("📝", "Summarizer", Inches(5.0), Inches(3.6)),
    ("🔍", "Fact\nExtractor", Inches(7.2), Inches(3.6)),
    ("🟦", "Debate\nLeft AI", Inches(9.4), Inches(3.6)),
    ("🟥", "Debate\nRight AI", Inches(11.4), Inches(3.6)),
]

for icon, name, x, y in agents:
    card = add_shape(slide, x, y, Inches(1.7), Inches(1.3), SURFACE, corner_radius=0.04)
    add_text(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.4),
             icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.6), Inches(1.5), Inches(0.6),
             name, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# More agents row 2
agents2 = [
    ("⚪", "Debate\nNeutral AI", Inches(1.7), Inches(5.2)),
    ("👤", "User Bias\nAnalyzer", Inches(3.9), Inches(5.2)),
    ("✅", "Quality\nGuard", Inches(6.1), Inches(5.2)),
    ("📅", "Timeline\nBuilder", Inches(8.3), Inches(5.2)),
    ("🌍", "Impact\nAnalyzer", Inches(10.5), Inches(5.2)),
]

for icon, name, x, y in agents2:
    card = add_shape(slide, x, y, Inches(1.7), Inches(1.3), SURFACE, corner_radius=0.04)
    add_text(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.4),
             icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.6), Inches(1.5), Inches(0.6),
             name, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Key points
add_shape(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.55), CARD_BG, corner_radius=0.1)
add_text(slide, Inches(0.8), Inches(6.73), Inches(11.7), Inches(0.5),
         "🤖 Ollama (Local LLM)  •  🔀 SelectorGroupChat (LLM picks next agent)  •  "
         "🔄 Retry & fallback  •  ✅ Every output validated by QualityGuard",
         font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 6: TECH STACK
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE_ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "⚙️ TECH STACK", font_size=32, color=BLUE_ACCENT, bold=True)

# Backend column
add_shape(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(5.5), SURFACE, corner_radius=0.05)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3.4), Inches(0.4),
         "🔧 Backend", font_size=20, color=GREEN, bold=True)
backend_items = [
    "Python 3.11+ (Runtime)",
    "FastAPI (Async Web Framework)",
    "AutoGen 0.4 (AI Orchestration)",
    "Ollama (Local LLM Inference)",
    "PostgreSQL 16 (Database)",
    "Redis 7 (Cache + Sessions)",
    "Celery (Task Queue)",
    "SQLAlchemy 2.0 (ORM)",
    "Stripe (Payments)",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(3.4), Inches(4.5),
                [f"• {item}" for item in backend_items], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# Frontend column
add_shape(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(5.5), SURFACE, corner_radius=0.05)
add_text(slide, Inches(5.0), Inches(1.5), Inches(3.4), Inches(0.4),
         "🎨 Frontend", font_size=20, color=PURPLE, bold=True)
frontend_items = [
    "Next.js 14 (App Router, SSR)",
    "TypeScript (Type Safety)",
    "Tailwind CSS + shadcn/ui",
    "Framer Motion (Animations)",
    "Zustand (State Mgmt)",
    "TanStack Query (Data Fetching)",
    "Recharts (Visualizations)",
    "Socket.IO (Real-time)",
    "Zod (Validation)",
]
add_bullet_list(slide, Inches(5.0), Inches(2.1), Inches(3.4), Inches(4.5),
                [f"• {item}" for item in frontend_items], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# Infrastructure column
add_shape(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), SURFACE, corner_radius=0.05)
add_text(slide, Inches(9.2), Inches(1.5), Inches(3.4), Inches(0.4),
         "🏗️ Infrastructure", font_size=20, color=AMBER, bold=True)
infra_items = [
    "Docker + Docker Compose",
    "Nginx / Traefik (LB + SSL)",
    "GitHub Actions (CI/CD)",
    "Prometheus + Grafana",
    "Structured JSON Logging",
    "JWT + OAuth 2.0 (Auth)",
    "Redis Rate Limiting",
    "OWASP Top 10 Compliant",
    "WCAG AA Accessible",
]
add_bullet_list(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(4.5),
                [f"• {item}" for item in infra_items], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# ═══════════════════════════════════════════════════════
# SLIDE 7: THEME ALIGNMENT
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), AMBER)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "🔥 THEME 1: REIMAGINING LEARNING & COLLABORATION", font_size=28, color=AMBER, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
         '"How can we reinvent the way people learn, share knowledge, and collaborate?"',
         font_size=20, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Alignment cards
alignments = [
    ("🎯", "Personalized Learning",
     "Bias Detection Engine creates a personalized\nlearning loop — tracks what you read, reveals\nyour blind spots, recommends what you're missing.",
     "Guiding: \"Make learning personalized, on-demand\""),
    ("🌍", "Collaborate Across Divides",
     "3-Side View bridges ideological \"regions\" —\nthe hardest kind of divide to collaborate across.\nHelps people understand viewpoints beyond their bubble.",
     "Guiding: \"Collaborate across roles, regions\""),
    ("💡", "Foster Curiosity",
     "AI Debate Mode is a curiosity engine — users\nengage with 3 AI personas arguing different\nperspectives. Gamified badges reward balanced reading.",
     "Guiding: \"Foster curiosity & knowledge-sharing\""),
    ("📊", "Information → Insight",
     "Raw biased news (information) is transformed by\nAI into balanced multi-perspective summaries\n(insight). Passive consumption → active thinking.",
     "Guiding: \"Digital tools turn information into insight\""),
]

for i, (icon, title, desc, guiding) in enumerate(alignments):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(2.3 + row * 2.5)

    card = add_shape(slide, x, y, Inches(5.9), Inches(2.2), SURFACE, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
             icon, font_size=28)
    add_text(slide, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.4),
             title, font_size=18, color=AMBER, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(5.3), Inches(1.0),
             desc, font_size=12, color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(1.7), Inches(5.3), Inches(0.3),
             guiding, font_size=10, color=GRAY)

# ═══════════════════════════════════════════════════════
# SLIDE 8: BUSINESS MODEL
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "💰 SUBSCRIPTION MODEL", font_size=32, color=GREEN, bold=True)

# Free tier
add_shape(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(5.5), SURFACE, corner_radius=0.05)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3.4), Inches(0.4),
         "FREE", font_size=24, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(2.0), Inches(3.4), Inches(0.5),
         "$0 / month", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
free_features = [
    "✅ 3-Side View (5/day)",
    "✅ AI Summaries (5/day)",
    "✅ AI Debate (2/day)",
    "✅ Basic Bias Score",
    "✅ 7-day Reading History",
    "❌ Timeline View",
    "❌ Impact Analyzer",
    "❌ Data Export",
    "⚠️ Minimal ads",
]
add_bullet_list(slide, Inches(0.9), Inches(2.7), Inches(3.2), Inches(3.8),
                free_features, font_size=12, color=LIGHT_GRAY, spacing=Pt(5))

# Premium tier
add_shape(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(5.5), RGBColor(0x1A, 0x2A, 0x44), corner_radius=0.05)
# Premium border glow effect
add_shape(slide, Inches(4.75), Inches(1.25), Inches(3.9), Inches(0.08), PRIMARY)
add_text(slide, Inches(5.0), Inches(1.4), Inches(3.4), Inches(0.3),
         "⭐ MOST POPULAR", font_size=10, color=AMBER, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(1.7), Inches(3.4), Inches(0.4),
         "PREMIUM", font_size=24, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(2.2), Inches(3.4), Inches(0.5),
         "$9.99 / month", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
premium_features = [
    "✅ Unlimited 3-Side Views",
    "✅ Unlimited AI Summaries",
    "✅ Unlimited AI Debates",
    "✅ Real-time Bias Score",
    "✅ Unlimited History",
    "✅ Timeline View",
    "✅ Impact Analyzer",
    "✅ PDF/CSV Export",
    "✅ No ads",
]
add_bullet_list(slide, Inches(5.1), Inches(2.9), Inches(3.2), Inches(3.8),
                premium_features, font_size=12, color=LIGHT_GRAY, spacing=Pt(5))

# Enterprise tier
add_shape(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), SURFACE, corner_radius=0.05)
add_text(slide, Inches(9.2), Inches(1.5), Inches(3.4), Inches(0.4),
         "ENTERPRISE", font_size=24, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.2), Inches(2.0), Inches(3.4), Inches(0.5),
         "$49.99 / month", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
enterprise_features = [
    "✅ Everything in Premium",
    "✅ API Access",
    "✅ Team Dashboard (50 users)",
    "✅ Custom News Sources",
    "✅ Custom AI Personas",
    "✅ Priority AI Processing",
    "✅ Priority Support (4h)",
    "✅ Slack Integration",
    "✅ SSO / SAML",
]
add_bullet_list(slide, Inches(9.3), Inches(2.7), Inches(3.2), Inches(3.8),
                enterprise_features, font_size=12, color=LIGHT_GRAY, spacing=Pt(5))

# ═══════════════════════════════════════════════════════
# SLIDE 9: IMPACT & WHY WE WIN
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), AMBER)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "🏆 WHY THIS IDEA WINS", font_size=32, color=AMBER, bold=True)

# Winning reasons
reasons = [
    ("✅", "Massive Social Impact", "Combats misinformation and polarization — a global crisis affecting billions.", GREEN),
    ("✅", "Real AI Usage", "AutoGen 0.4 multi-agent system with 12 specialized agents. Not a GPT wrapper.", PRIMARY),
    ("✅", "Stunning Visual Demo", "3-Side View is immediately impressive. Judges understand it in seconds.", PURPLE),
    ("✅", "Truly Unique", "Not another chatbot, fitness app, or to-do list. Solves a real cognitive problem.", AMBER),
    ("✅", "Geopolitically Relevant", "Misinformation, elections, wars — this is the most timely problem to solve.", RED_ACCENT),
    ("✅", "Production Architecture", "Scalable, secure, containerized. Not a hackathon throwaway — a real product.", BLUE_ACCENT),
]

for i, (icon, title, desc, accent) in enumerate(reasons):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.3 + row * 1.9)
    card = add_shape(slide, x, y, Inches(5.9), Inches(1.6), SURFACE, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.4),
             f"{icon} {title}", font_size=18, color=accent, bold=True)
    add_text(slide, x + Inches(0.6), y + Inches(0.7), Inches(5.0), Inches(0.7),
             desc, font_size=13, color=LIGHT_GRAY)

# Bottom quote
add_shape(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.7), CARD_BG, corner_radius=0.1)
add_text(slide, Inches(2.7), Inches(6.55), Inches(7.9), Inches(0.6),
         "\"We are not building another news app. We are building a system that helps people think, not just react.\"",
         font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 10: THANK YOU / CTA
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), PRIMARY)

# Coin icon
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.0), Inches(1.5), Inches(1.5))
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY
circle.line.fill.background()
add_text(slide, Inches(6.05), Inches(1.35), Inches(1.2), Inches(0.8), "⚖️",
         font_size=44, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(2.8), Inches(8.3), Inches(0.8),
         "BothSidesOfACoin", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.6),
         "Breaking the echo chamber by showing all sides of reality.",
         font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(4.5), Inches(8.3), Inches(0.6),
         "Thank You", font_size=36, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Links
add_shape(slide, Inches(3.8), Inches(5.5), Inches(5.7), Inches(1.2), SURFACE, corner_radius=0.08)
add_text(slide, Inches(4.0), Inches(5.6), Inches(5.3), Inches(0.35),
         "🔗 github.com/ShubhamBrody/BothSidesOfACoin",
         font_size=14, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.0), Inches(5.95), Inches(5.3), Inches(0.35),
         "🔥 Theme 1: Reimagining Learning & Collaboration",
         font_size=13, color=AMBER, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.0), Inches(6.3), Inches(5.3), Inches(0.35),
         "📧 Team BothSidesOfACoin  •  April 2026",
         font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = r"d:\Study Material\AI Projects\BothSidesOfACoin\docs\BothSidesOfACoin_Ideation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
